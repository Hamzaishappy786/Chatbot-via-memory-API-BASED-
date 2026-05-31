import base64
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from app.ingestion.parsers.base import BaseParser, ParsedContent


class PPTXParser(BaseParser):
    @staticmethod
    def supported_extensions() -> list[str]:
        return [".pptx"]

    def parse(self, file_path: str) -> ParsedContent:
        prs = Presentation(file_path)
        content = ParsedContent()

        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_bytes = shape.image.blob
                        img_b64 = base64.b64encode(img_bytes).decode()
                        content.images.append({
                            "image_base64": img_b64,
                            "page": slide_num,
                            "source": "slide_image",
                        })
                    except Exception:
                        continue

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        header = rows[0]
                        separator = " | ".join(["---"] * len(table.columns))
                        table_md = "\n".join([header, separator] + rows[1:])
                        texts.append(f"\n{table_md}\n")

            if texts:
                content.text_blocks.append({
                    "text": "\n".join(texts),
                    "page": slide_num,
                })

        return content
